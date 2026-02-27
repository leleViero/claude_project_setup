from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT = r"C:\Users\DAVI\Documents\02 LocalRepo\claude_project_setup\test-presentation.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

DARK_BLUE  = RGBColor(0x1C, 0x28, 0x33)
MID_BLUE   = RGBColor(0x2E, 0x40, 0x53)
SILVER     = RGBColor(0xAA, 0xB7, 0xB8)
OFF_WHITE  = RGBColor(0xF4, 0xF6, 0xF6)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, l, t, w, h, text, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


# ── Slide 1: Title ────────────────────────────────────────────────────────────
sl1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(sl1, DARK_BLUE)

# Left accent bar
add_rect(sl1, 0, 0, 0.35, 7.5, MID_BLUE)

# Decorative bottom band
add_rect(sl1, 0.35, 6.2, 12.98, 1.3, MID_BLUE)

# Title
add_text_box(sl1, 1.0, 2.2, 11.0, 1.5,
             "Test Presentation",
             54, bold=True, color=OFF_WHITE, align=PP_ALIGN.LEFT)

# Subtitle
add_text_box(sl1, 1.0, 3.9, 11.0, 0.9,
             "A Simple Three-Slide Demo",
             24, bold=False, color=SILVER, align=PP_ALIGN.LEFT)

# Date
add_text_box(sl1, 1.0, 6.35, 6.0, 0.6,
             "February 2026",
             14, bold=False, color=SILVER, align=PP_ALIGN.LEFT)


# ── Slide 2: Content / Bullet Points ──────────────────────────────────────────
sl2 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(sl2, OFF_WHITE)

# Top header bar
add_rect(sl2, 0, 0, 13.33, 1.4, DARK_BLUE)

# Left accent
add_rect(sl2, 0, 1.4, 0.35, 6.1, MID_BLUE)

# Slide title
add_text_box(sl2, 0.55, 0.25, 12.0, 0.9,
             "Key Points",
             36, bold=True, color=OFF_WHITE)

# Bullet points — added one at a time for clean control
bullets = [
    "This is the first bullet point — introduce your main idea here.",
    "This is the second bullet point — elaborate with supporting details.",
    "This is the third bullet point — provide evidence or examples.",
    "This is the fourth bullet point — highlight any key takeaways.",
    "This is the fifth bullet point — wrap up the section's argument.",
]

for i, text in enumerate(bullets):
    top = 1.7 + i * 0.9
    # Bullet marker
    add_rect(sl2, 0.65, top + 0.22, 0.18, 0.18, MID_BLUE)
    # Bullet text
    add_text_box(sl2, 1.05, top, 11.8, 0.8,
                 text, 18, bold=False, color=DARK_BLUE)


# ── Slide 3: Closing / Thank You ──────────────────────────────────────────────
sl3 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(sl3, DARK_BLUE)

# Centre decorative rectangle
add_rect(sl3, 3.0, 2.5, 7.33, 2.5, MID_BLUE)

# "Thank You" heading
add_text_box(sl3, 3.0, 2.7, 7.33, 1.1,
             "Thank You",
             52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Closing line
add_text_box(sl3, 3.0, 3.9, 7.33, 0.8,
             "Questions? Reach out any time.",
             20, bold=False, color=SILVER, align=PP_ALIGN.CENTER)

# Bottom note
add_text_box(sl3, 0, 6.8, 13.33, 0.5,
             "Test Presentation  |  February 2026",
             12, bold=False, color=SILVER, align=PP_ALIGN.CENTER)


prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
