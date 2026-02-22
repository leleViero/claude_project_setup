from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Color Palette (from reference) ---
DARK_NAVY = RGBColor(0x0D, 0x06, 0x26)
MUTED_PURPLE = RGBColor(0x61, 0x5D, 0x71)
CORAL = RGBColor(0xE2, 0x5B, 0x67)
SALMON = RGBColor(0xF6, 0x76, 0x5D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF7, 0xF6, 0xF9)
ACCENT_TEAL = RGBColor(0x3A, 0x7C, 0x89)

FONT = 'Arial'
W = Inches(16)
H = Inches(9)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=20, color=DARK_NAVY, bold=False, alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return txBox


def add_header_bar(slide, title_text):
    """Add the signature dark header bar with white title, matching the reference."""
    bar = add_rect(slide, Inches(0), Inches(0), W, Inches(1.15), DARK_NAVY)
    add_text(slide, Inches(0.8), Inches(0.22), Inches(14), Inches(0.8),
             title_text, font_size=36, color=WHITE, bold=True)
    return bar


def add_slide():
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    add_bg(slide, LIGHT_BG)
    return slide


# ============================================================
# SLIDE 0: TITLE - "HYPERVOLUME"
# ============================================================
slide = add_slide()
add_bg(slide, DARK_NAVY)

# Large brand name
add_text(slide, Inches(0), Inches(2.4), W, Inches(1.5),
         'HYPERVOLUME', font_size=72, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Tagline
add_text(slide, Inches(0), Inches(3.9), W, Inches(0.8),
         'The Fabric of Trust', font_size=44, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)

# Subtle accent line
add_rect(slide, Inches(6.5), Inches(5.0), Inches(3), Inches(0.06), CORAL)

# Subtitle
add_text(slide, Inches(0), Inches(5.5), W, Inches(0.7),
         'Trust, Transparency & Transferability as a Service', font_size=22, color=MUTED_PURPLE, bold=False, alignment=PP_ALIGN.CENTER)

# Bottom text
add_text(slide, Inches(0), Inches(7.8), W, Inches(0.5),
         '90-Second Investor Pitch', font_size=16, color=MUTED_PURPLE, bold=False, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 1: THE PROBLEM
# ============================================================
slide = add_slide()
add_header_bar(slide, 'The Problem')

# Story block - left side
add_text(slide, Inches(0.8), Inches(1.6), Inches(7), Inches(1.2),
         'A manufacturer fails a compliance audit.',
         font_size=30, color=DARK_NAVY, bold=True)

add_text(slide, Inches(0.8), Inches(2.7), Inches(7), Inches(1.0),
         'Not because they did anything wrong --\nbecause they cannot prove they didn\'t.',
         font_size=22, color=MUTED_PURPLE, bold=False)

# Right side - the pain points in accent boxes
box_x = Inches(8.8)
box_w = Inches(6.2)

# Pain point cards
cards = [
    ('3 ERPs, 2 countries, 4 vendors', 'Data scattered. Nothing verifiable.'),
    ('Lost contracts', 'Legal exposure, reputational damage.'),
    ('Not an edge case', 'This is the default state of enterprise data.'),
]

for idx, (title, subtitle) in enumerate(cards):
    y = Inches(1.6) + Inches(idx * 2.1)
    card = add_rect(slide, box_x, y, box_w, Inches(1.7), WHITE)
    # left accent bar on card
    add_rect(slide, box_x, y, Inches(0.08), Inches(1.7), CORAL)
    add_text(slide, box_x + Inches(0.35), y + Inches(0.25), box_w - Inches(0.6), Inches(0.6),
             title, font_size=22, color=DARK_NAVY, bold=True)
    add_text(slide, box_x + Inches(0.35), y + Inches(0.9), box_w - Inches(0.6), Inches(0.6),
             subtitle, font_size=18, color=MUTED_PURPLE)

# Bottom sector tags
sectors = ['Healthcare', 'Logistics', 'ESG Reporting', 'Supply Chain']
for idx, sector in enumerate(sectors):
    x = Inches(0.8) + Inches(idx * 2.2)
    pill = add_rect(slide, x, Inches(7.4), Inches(1.9), Inches(0.5), DARK_NAVY)
    pill.text_frame.word_wrap = False
    p = pill.text_frame.paragraphs[0]
    p.text = sector
    p.font.name = FONT
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    pill.text_frame.paragraphs[0].space_before = Pt(4)


# ============================================================
# SLIDE 2: THE SHIFT
# ============================================================
slide = add_slide()
add_header_bar(slide, 'The Shift')

# Old paradigm
add_text(slide, Inches(0.8), Inches(1.6), Inches(6.5), Inches(0.6),
         'The last 20 years:', font_size=20, color=MUTED_PURPLE)

old_words = ['Velocity', 'Variety', 'Volume']
for idx, word in enumerate(old_words):
    x = Inches(0.8) + Inches(idx * 2.5)
    box = add_rect(slide, x, Inches(2.2), Inches(2.2), Inches(1.4), WHITE)
    add_text(slide, x, Inches(2.45), Inches(2.2), Inches(1.0),
             word, font_size=28, color=MUTED_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

# Arrow / transition
add_text(slide, Inches(0.8), Inches(4.0), Inches(7), Inches(0.6),
         'The next decade demands:', font_size=20, color=DARK_NAVY, bold=True)

new_words = ['Privacy', 'Transferability', 'Transparency']
for idx, word in enumerate(new_words):
    x = Inches(0.8) + Inches(idx * 2.5)
    box = add_rect(slide, x, Inches(4.6), Inches(2.2), Inches(1.4), CORAL)
    add_text(slide, x, Inches(4.85), Inches(2.2), Inches(1.0),
             word, font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Right side - the key insight
insight_box = add_rect(slide, Inches(8.8), Inches(1.6), Inches(6.2), Inches(5.5), WHITE)
add_rect(slide, Inches(8.8), Inches(1.6), Inches(6.2), Inches(0.08), CORAL)

add_text(slide, Inches(9.2), Inches(2.2), Inches(5.4), Inches(1.2),
         'Existing solutions always\nsacrifice at least one.',
         font_size=30, color=DARK_NAVY, bold=True)

add_text(slide, Inches(9.2), Inches(3.8), Inches(5.4), Inches(0.8),
         'Always a trade-off.', font_size=26, color=CORAL, bold=True)

add_text(slide, Inches(9.2), Inches(5.2), Inches(5.4), Inches(1.2),
         'Hypervolume eliminates\nthat trade-off.',
         font_size=30, color=DARK_NAVY, bold=True)


# ============================================================
# SLIDE 3: THE SOLUTION - T3aaS
# ============================================================
slide = add_slide()
add_header_bar(slide, 'Hypervolume T3aaS')

# Main value prop
add_text(slide, Inches(0.8), Inches(1.6), Inches(14), Inches(1.0),
         'Trust, Transparency & Transferability as a Service',
         font_size=34, color=DARK_NAVY, bold=True)

add_text(slide, Inches(0.8), Inches(2.5), Inches(14), Inches(0.7),
         'A trust infrastructure layer. Not middleware. Not blockchain. Infrastructure.',
         font_size=22, color=MUTED_PURPLE)

# Three pillars
pillars = [
    ('Proprietary\nPersistence Layer', 'Novel database paradigm\nfor immutable records'),
    ('Consensus\nMechanism', 'Cross-boundary verification\nwithout exposing data'),
    ('Universal\nAPI Surface', 'Connects to any system.\nNo forced migration.'),
]

for idx, (title, desc) in enumerate(pillars):
    x = Inches(0.8) + Inches(idx * 5.0)
    # Card
    card = add_rect(slide, x, Inches(3.6), Inches(4.5), Inches(3.8), WHITE)
    # Top accent
    add_rect(slide, x, Inches(3.6), Inches(4.5), Inches(0.08), CORAL if idx == 0 else DARK_NAVY if idx == 1 else ACCENT_TEAL)
    # Number circle
    num_color = CORAL if idx == 0 else DARK_NAVY if idx == 1 else ACCENT_TEAL
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), Inches(4.0), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = num_color
    circle.line.fill.background()
    p = circle.text_frame.paragraphs[0]
    p.text = str(idx + 1)
    p.font.name = FONT
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text(slide, x + Inches(0.3), Inches(4.9), Inches(3.8), Inches(1.2),
             title, font_size=24, color=DARK_NAVY, bold=True)
    add_text(slide, x + Inches(0.3), Inches(6.1), Inches(3.8), Inches(1.0),
             desc, font_size=18, color=MUTED_PURPLE)

# Bottom tagline
add_text(slide, Inches(0.8), Inches(7.8), Inches(14), Inches(0.6),
         'Native Selective Disclosure  |  Universal Provenance Protocol  |  Domain-Aware Language',
         font_size=16, color=MUTED_PURPLE, alignment=PP_ALIGN.LEFT)


# ============================================================
# SLIDE 4: THE MARKET
# ============================================================
slide = add_slide()
add_header_bar(slide, 'The Market')

# Giant number - left
add_text(slide, Inches(0.8), Inches(1.5), Inches(7), Inches(1.8),
         '€900B+', font_size=96, color=CORAL, bold=True)

add_text(slide, Inches(0.8), Inches(3.3), Inches(7), Inches(0.8),
         'Annual EU Compliance Burden', font_size=28, color=DARK_NAVY, bold=True)

# The gap stat
add_text(slide, Inches(0.8), Inches(4.5), Inches(7), Inches(1.8),
         '<3%', font_size=96, color=DARK_NAVY, bold=True)

add_text(slide, Inches(0.8), Inches(6.3), Inches(7), Inches(0.8),
         'Addressed by technology today', font_size=28, color=MUTED_PURPLE, bold=True)

# Right side - market funnel
funnel_x = Inches(8.8)
funnel_w = Inches(6.2)

funnel_data = [
    ('€900B+', 'EU Annual Compliance', Inches(6.2), DARK_NAVY),
    ('€250B', 'Target Segment (TAM)', Inches(5.2), MUTED_PURPLE),
    ('€12.5B', 'Serviceable Market (SAM)', Inches(4.2), CORAL),
    ('€125M', 'Target 1% (SOM)', Inches(3.2), SALMON),
]

for idx, (amount, label, width, color) in enumerate(funnel_data):
    y = Inches(1.8) + Inches(idx * 1.6)
    offset = (funnel_w - width) / 2
    bar = add_rect(slide, funnel_x + offset, y, width, Inches(1.1), color)
    # Amount text
    p = bar.text_frame.paragraphs[0]
    p.text = f'{amount}  {label}'
    p.font.name = FONT
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    bar.text_frame.paragraphs[0].space_before = Pt(8)

add_text(slide, Inches(8.8), Inches(8.0), Inches(6.2), Inches(0.5),
         'MARKET FUNNEL', font_size=14, color=MUTED_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5: THE ASK
# ============================================================
slide = add_slide()
add_bg(slide, DARK_NAVY)

# Title
add_text(slide, Inches(0), Inches(1.0), W, Inches(1.0),
         'The Ask', font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Accent line
add_rect(slide, Inches(6.5), Inches(2.2), Inches(3), Inches(0.06), CORAL)

# Amount
add_text(slide, Inches(0), Inches(2.8), W, Inches(1.5),
         '€1M', font_size=96, color=CORAL, bold=True, alignment=PP_ALIGN.CENTER)

# Purpose
add_text(slide, Inches(0), Inches(4.5), W, Inches(0.8),
         'First enterprise pilots: Supply Chain & Healthcare',
         font_size=28, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)

# Key message boxes
msgs = [
    'The whitepaper exists.',
    'The codebase exists.',
    'The market gap exists.',
]
for idx, msg in enumerate(msgs):
    x = Inches(2.5) + Inches(idx * 3.8)
    card = add_rect(slide, x, Inches(5.8), Inches(3.3), Inches(1.0), RGBColor(0x1A, 0x10, 0x35))
    add_text(slide, x, Inches(5.95), Inches(3.3), Inches(0.8),
             msg, font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Bottom CTA
add_text(slide, Inches(0), Inches(7.4), W, Inches(0.8),
         'Trust infrastructure is not a feature -- it is the foundation.',
         font_size=24, color=MUTED_PURPLE, bold=False, alignment=PP_ALIGN.CENTER)

# Contact
add_text(slide, Inches(0), Inches(8.2), W, Inches(0.5),
         'info@hypervolume.io', font_size=16, color=MUTED_PURPLE, alignment=PP_ALIGN.CENTER)


# ============================================================
# SAVE
# ============================================================
output_path = r'C:\Users\DAVI\Documents\02 LocalRepo\claude_project_setup\Hypervolume-Investor-Pitch-90s.pptx'
prs.save(output_path)
print(f'Saved: {output_path}')
